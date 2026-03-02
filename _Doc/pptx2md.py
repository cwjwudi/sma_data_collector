#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""将 .pptx 转为 Markdown：按幻灯片提取标题与正文。"""
import sys
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("请先安装: pip install python-pptx")
    sys.exit(1)


def extract_text(shape):
    """从 shape 提取纯文本。"""
    if not hasattr(shape, "text_frame"):
        return ""
    text = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            text.append(line)
    return "\n".join(text).strip()


def pptx_to_markdown(pptx_path: Path, md_path: Path = None) -> str:
    prs = Presentation(str(pptx_path))
    lines = [f"# {pptx_path.stem}\n"]

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n---\n\n## 幻灯片 {i}\n")
        for shape in slide.shapes:
            text = extract_text(shape)
            if not text:
                continue
            # 若整段较短或像标题（无句号等），当作标题
            if len(text) < 80 and not re.search(r"[。.!?；;]\s*$", text):
                lines.append(f"### {text}\n")
            else:
                for paragraph in text.split("\n"):
                    if paragraph.strip():
                        lines.append(f"{paragraph.strip()}\n")
                lines.append("")

    md_content = "\n".join(lines).strip()
    if md_path is None:
        md_path = pptx_path.with_suffix(".md")
    md_path.write_text(md_content, encoding="utf-8")
    return str(md_path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python pptx2md.py <文件.pptx> [输出.md]")
        sys.exit(1)
    ppt = Path(sys.argv[1]).resolve()
    out = Path(sys.argv[2]).resolve() if len(sys.argv) > 2 else None
    if not ppt.exists():
        print(f"文件不存在: {ppt}")
        sys.exit(1)
    result = pptx_to_markdown(pptx_path=ppt, md_path=out)
    print(f"已生成: {result}")
